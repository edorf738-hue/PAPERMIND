import os
import re
import uuid
import json
import io
import logging
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from dotenv import load_dotenv
import fitz  # pymupdf
import docx
from pptx import Presentation
import openpyxl
import xlrd
from PIL import Image
import pytesseract
from supabase import create_client, Client
from groq import Groq


load_dotenv()

logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
logger = logging.getLogger(__name__)

_tesseract_path = os.getenv("TESSERACT_PATH", r"C:\Program Files\Tesseract-OCR\tesseract.exe")
pytesseract.pytesseract.tesseract_cmd = _tesseract_path
if os.path.exists(_tesseract_path):
    logger.info(f"Tesseract ditemukan: {_tesseract_path}")
else:
    logger.warning(f"Tesseract TIDAK ditemukan di: {_tesseract_path} — OCR tidak akan berfungsi")

app = FastAPI(title="PaperMind API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 
supabase: Client = create_client(
    os.getenv("SUPABASE_URL"),
    os.getenv("SUPABASE_SERVICE_KEY"),
)


groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))

ANALYZE_MODEL = "llama-3.3-70b-versatile"
CHAT_MODEL = "llama-3.3-70b-versatile"
FALLBACK_MODEL = "llama-3.1-8b-instant"


def groq_create(messages: list, max_tokens: int, temperature: float = 0.4, response_format: dict | None = None, model: str = CHAT_MODEL, allow_fallback: bool = True):
    """Panggil Groq, dengan fallback opsional ke FALLBACK_MODEL jika rate limit."""
    kwargs = dict(model=model, messages=messages, temperature=temperature, max_tokens=max_tokens)
    if response_format:
        kwargs["response_format"] = response_format
    try:
        return groq_client.chat.completions.create(**kwargs)
    except Exception as e:
        err = str(e)
        if "rate_limit_exceeded" in err and model != FALLBACK_MODEL and allow_fallback:
            logger.warning(f"Rate limit {model} — fallback ke {FALLBACK_MODEL} dengan konteks ringkas")
            sys_msg = next((m for m in messages if m["role"] == "system"), None)
            user_msgs = [m for m in messages if m["role"] == "user"]
            fallback_msgs = []
            if sys_msg:
                short = sys_msg["content"][:3500]
                if len(sys_msg["content"]) > 3500:
                    short += "\n...[konteks dipotong karena keterbatasan kuota]..."
                fallback_msgs.append({"role": "system", "content": short})
            if user_msgs:
                last_user = user_msgs[-1]
                if len(last_user["content"]) > 3500:
                    last_user = {"role": "user", "content": last_user["content"][:3500] + "\n...[dipotong]..."}
                fallback_msgs.append(last_user)
            kwargs["model"] = FALLBACK_MODEL
            kwargs["messages"] = fallback_msgs
            kwargs["max_tokens"] = min(max_tokens, 500)
            kwargs.pop("response_format", None)
            return groq_client.chat.completions.create(**kwargs)
        raise


ALLOWED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".txt", ".md"
}

EXTENSION_LABELS = {
    ".pdf": "PDF",
    ".docx": "Word",
    ".doc": "Word",
    ".pptx": "PowerPoint",
    ".ppt": "PowerPoint",
    ".xlsx": "Excel",
    ".xls": "Excel",
    ".txt": "Teks",
    ".md": "Markdown",
}


def extract_text(file_bytes: bytes, filename: str) -> tuple[str, int]:
    ext = os.path.splitext(filename.lower())[1]

    if ext == ".pdf":
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        page_count = len(doc)
        parts = []
        ocr_used = False

        for i, page in enumerate(doc):
            # Method 1: standard text
            text = page.get_text().strip()

            # Method 2: per-block (lebih baik untuk PDF dari PPT/Keynote)
            if not text:
                blocks = page.get_text("blocks")
                text = "\n".join(
                    b[4].strip() for b in blocks
                    if len(b) > 4 and isinstance(b[4], str) and b[4].strip()
                )

            # Method 3: per-span dari dict
            if not text:
                raw = page.get_text("dict")
                spans = []
                for block in raw.get("blocks", []):
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            t = span.get("text", "").strip()
                            if t:
                                spans.append(t)
                text = " ".join(spans)

            # Method 4: OCR fallback untuk halaman yang benar-benar image-based
            if not text:
                logger.info(f"Halaman {i+1}: teks kosong, mencoba OCR...")
                try:
                    mat = fitz.Matrix(2.0, 2.0)
                    pix = page.get_pixmap(matrix=mat)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    try:
                        text = pytesseract.image_to_string(img, lang="ind+eng").strip()
                        logger.info(f"Halaman {i+1}: OCR ind+eng berhasil ({len(text)} karakter)")
                    except Exception as ocr_err:
                        logger.warning(f"Halaman {i+1}: OCR ind+eng gagal ({ocr_err}), mencoba eng saja...")
                        text = pytesseract.image_to_string(img, lang="eng").strip()
                        logger.info(f"Halaman {i+1}: OCR eng berhasil ({len(text)} karakter)")
                    if text:
                        ocr_used = True
                except Exception as e:
                    logger.error(f"Halaman {i+1}: OCR gagal total — {e}")

            if text:
                parts.append(f"[Page {i + 1}]\n{text}")

        doc.close()
        return "\n\n".join(parts), page_count

    if ext in (".docx", ".doc"):
        try:
            doc = docx.Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            page_count = max(1, len(paragraphs) // 25)
            return "\n\n".join(paragraphs), page_count
        except Exception as e:
            if "zip" in str(e).lower():
                raise ValueError("File .doc format lama (Word 97-2003) tidak didukung. Silakan buka di Microsoft Word atau Google Docs lalu simpan ulang sebagai .docx.")
            raise

    if ext in (".pptx", ".ppt"):
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            parts = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    parts.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
            return "\n\n".join(parts), len(prs.slides)
        except Exception as e:
            if "zip" in str(e).lower():
                raise ValueError("File .ppt format lama (PowerPoint 97-2003) tidak didukung. Silakan buka di PowerPoint atau Google Slides lalu simpan ulang sebagai .pptx.")
            raise

    if ext in (".xlsx", ".xls"):
        if ext == ".xls":
            try:
                wb = xlrd.open_workbook(file_contents=file_bytes)
                parts = []
                for sheet in wb.sheets():
                    rows = []
                    for row_idx in range(sheet.nrows):
                        cells = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols) if sheet.cell_value(row_idx, col_idx) != ""]
                        if cells:
                            rows.append("\t".join(cells))
                    if rows:
                        parts.append(f"[Sheet: {sheet.name}]\n" + "\n".join(rows))
                return "\n\n".join(parts), wb.nsheets
            except Exception as e:
                raise ValueError(f"Gagal membaca file .xls: {str(e)}")
        else:
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        rows.append("\t".join(cells))
                if rows:
                    parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
            return "\n\n".join(parts), len(wb.worksheets)

    if ext in (".txt", ".md"):
        text = file_bytes.decode("utf-8", errors="ignore")
        lines = [l for l in text.splitlines() if l.strip()]
        return "\n".join(lines), max(1, len(lines) // 40)

    raise ValueError(f"Format file tidak didukung: {ext}")


def parse_rate_limit_error(err: str) -> str:
    limit_match = re.search(r'Limit (\d+)', err)
    used_match = re.search(r'Used (\d+)', err)
    retry_match = re.search(r'try again in ([\dh\sm.]+?)(?:\.?\s*Need|$)', err)

    limit = int(limit_match.group(1)) if limit_match else None
    used = int(used_match.group(1)) if used_match else None

    retry_str = ""
    if retry_match:
        raw_time = retry_match.group(1).strip()
        h = re.search(r'(\d+)h', raw_time)
        m = re.search(r'(\d+)m', raw_time)
        hours = int(h.group(1)) if h else 0
        minutes = int(m.group(1)) if m else 0
        parts = []
        if hours:
            parts.append(f"{hours} jam")
        if minutes:
            parts.append(f"{minutes} menit")
        retry_str = "~" + " ".join(parts) if parts else "beberapa menit"

    lines = ["Kuota Groq API hari ini sudah habis."]
    if limit:
        lines.append(f"- Limit: {limit:,} token/hari".replace(",", "."))
    if used:
        lines.append(f"- Sudah dipakai: {used:,} token".replace(",", "."))
    if retry_str:
        lines.append(f"- Coba lagi dalam: {retry_str}")

    return "\n".join(lines)


def parse_json_from_ai(raw: str) -> dict:
    """Ekstrak JSON dari respons AI, toleran terhadap teks sebelum/sesudah JSON."""
    cleaned = raw.strip()

    # Hapus markdown code block
    code_match = re.search(r'```(?:json)?\s*([\s\S]*?)```', cleaned)
    if code_match:
        cleaned = code_match.group(1).strip()

    # Coba parse langsung
    try:
        return json.loads(cleaned)
    except (json.JSONDecodeError, ValueError):
        pass

    # Coba temukan JSON object: dari { pertama sampai } terakhir
    start = cleaned.find('{')
    end = cleaned.rfind('}')
    if start != -1 and end > start:
        try:
            return json.loads(cleaned[start:end + 1])
        except (json.JSONDecodeError, ValueError):
            pass

    # Fallback: regex greedy untuk nested JSON
    brace_match = re.search(r'\{[\s\S]*\}', cleaned)
    if brace_match:
        try:
            return json.loads(brace_match.group())
        except (json.JSONDecodeError, ValueError):
            pass

    # Last resort: JSON dengan unescaped quotes — ekstrak field "output" langsung
    output_match = re.search(r'"output"\s*:\s*"([\s\S]*)"\s*,\s*"citations"', cleaned)
    if output_match:
        logger.warning("parse_json_from_ai: pakai regex fallback (JSON malformed, unescaped quotes)")
        return {"output": output_match.group(1), "citations": [], "quiz": None}

    logger.warning(f"parse_json_from_ai gagal. Raw ({len(raw)} chars): {raw[:400]}")
    raise ValueError("Tidak ditemukan JSON valid dalam respons AI")


def analyze_document(text: str, filename: str) -> dict:
    truncated = text[:8000]
    prompt = f"""You must respond with valid JSON only. No explanation, no markdown, no text before or after the JSON object.

Analyze the following document and extract its information.

File name: {filename}

Document text (partial):
{truncated}

Return ONLY this JSON object, nothing else:
{{
  "name": "judul lengkap dokumen",
  "topic": "topik utama dalam 1 kalimat pendek",
  "keywords": ["keyword1", "keyword2", "keyword3"],
  "language": "Bahasa Indonesia atau English",
  "author": "nama penulis atau Tidak ditemukan",
  "year": "tahun publikasi atau Tidak ditemukan",
  "type": "Jurnal Ilmiah atau Thesis atau Skripsi atau Laporan atau Makalah atau Presentasi atau Spreadsheet atau Dokumen",
  "suggestedQuestions": [
    "Apa tujuan utama penelitian ini?",
    "Metodologi apa yang digunakan?",
    "Apa kesimpulan utama penelitian?",
    "Apa keterbatasan penelitian ini?",
    "Apa research gap yang ditemukan?",
    "Buatkan 5 soal kuis dari dokumen ini"
  ]
}}"""

    try:
        response = groq_client.chat.completions.create(
            model=ANALYZE_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2,
            max_tokens=1000,
            response_format={"type": "json_object"},
        )
    except Exception as e:
        err = str(e)
        if "rate_limit_exceeded" in err:
            raise ValueError(parse_rate_limit_error(err))
        raise

    raw = response.choices[0].message.content.strip()
    logger.info(f"analyze_document raw ({len(raw)} chars): {raw[:200]}")

    return parse_json_from_ai(raw)


@app.get("/health")
def health():
    return {"status": "ok"}


@app.delete("/session/{session_id}")
async def delete_session(session_id: str):
    try:
        supabase.table("chat_history").delete().eq("session_id", session_id).execute()
        supabase.table("sessions").delete().eq("session_id", session_id).execute()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gagal menghapus sesi: {str(e)}")
    return {"success": True}


@app.get("/sessions")
async def get_user_sessions(user_id: str):
    try:
        result = supabase.table("sessions") \
            .select("session_id, metadata, created_at") \
            .eq("user_id", user_id) \
            .order("created_at", desc=True) \
            .limit(50) \
            .execute()
        return {"sessions": result.data or []}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gagal mengambil sesi: {str(e)}")


@app.post("/upload")
async def upload_document(file: UploadFile = File(...), user_id: str | None = Form(None)):
    ext = os.path.splitext(file.filename.lower())[1]
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Format tidak didukung. Format yang diterima: PDF, Word, PowerPoint, Excel, TXT."
        )

    file_bytes = await file.read()
    file_size_mb = len(file_bytes) / (1024 * 1024)

    if file_size_mb > 100:
        raise HTTPException(status_code=400, detail="Ukuran file melebihi batas 100MB.")

    try:
        doc_text, page_count = extract_text(file_bytes, file.filename)
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Gagal membaca file: {str(e)}")

    if len(doc_text.strip()) < 30:
        raise HTTPException(
            status_code=422,
            detail="Teks tidak berhasil diekstrak dari file ini. Kemungkinan PDF berbasis gambar (scan/sertifikat) dan OCR tidak dapat membaca teksnya. Pastikan Tesseract sudah terinstall, atau coba file dengan kualitas gambar lebih tinggi."
        )

    try:
        analysis = analyze_document(doc_text, file.filename)
    except json.JSONDecodeError:
        raise HTTPException(status_code=500, detail="AI gagal menganalisis struktur dokumen.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gagal menganalisis dokumen: {str(e)}")

    session_id = str(uuid.uuid4())

    metadata = {
        "name": analysis.get("name", file.filename.replace(".pdf", "")),
        "pages": page_count,
        "size": f"{file_size_mb:.1f} MB",
        "topic": analysis.get("topic", ""),
        "keywords": analysis.get("keywords", []),
        "language": analysis.get("language", ""),
        "author": analysis.get("author", ""),
        "year": analysis.get("year", ""),
        "type": analysis.get("type", ""),
    }

    try:
        supabase.table("sessions").insert({
            "session_id": session_id,
            "document_text": doc_text,
            "metadata": metadata,
            "user_id": user_id or None,
        }).execute()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gagal menyimpan sesi ke database: {str(e)}")

    return {
        "sessionId": session_id,
        "document": metadata,
        "suggestedQuestions": analysis.get("suggestedQuestions", []),
    }


@app.get("/session/{session_id}")
async def get_session(session_id: str):
    try:
        result = supabase.table("sessions") \
            .select("metadata") \
            .eq("session_id", session_id) \
            .single() \
            .execute()
        if not result.data:
            raise HTTPException(status_code=404, detail="Sesi tidak ditemukan.")
        metadata = result.data["metadata"]
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gagal mengambil sesi: {str(e)}")

    try:
        history_result = supabase.table("chat_history") \
            .select("role, content") \
            .eq("session_id", session_id) \
            .order("created_at") \
            .execute()
        raw_messages = history_result.data or []
    except Exception:
        raw_messages = []

    def clean_content(role, content):
        """Bersihkan content yang tersimpan sebagai raw JSON (akibat parse gagal di masa lalu)."""
        if role != "assistant" or not content:
            return content
        stripped = content.strip()
        if stripped.startswith('{'):
            try:
                parsed = json.loads(stripped)
                if isinstance(parsed, dict) and "output" in parsed:
                    return str(parsed["output"])
            except (json.JSONDecodeError, ValueError):
                pass
        return content

    messages = [
        {
            "id": str(i),
            "role": h["role"] if h["role"] in ("user", "quiz_result") else "ai",
            "content": clean_content(h["role"], h["content"]),
            "citations": [],
        }
        for i, h in enumerate(raw_messages)
        if not (h["role"] == "assistant" and h["content"] == "quiz")
    ]

    return {"document": metadata, "messages": messages}


class ChatRequest(BaseModel):
    sessionId: str
    chatInput: str
    action: str = "sendMessage"
    userId: str | None = None
    isQuiz: bool = False
    quizType: str = "pilgan"



@app.post("/chat")
async def chat(body: ChatRequest):
    try:
        result = supabase.table("sessions") \
            .select("document_text, metadata") \
            .eq("session_id", body.sessionId) \
            .single() \
            .execute()
        if not result.data:
            raise HTTPException(status_code=404, detail="Sesi tidak ditemukan.")
        doc_text = result.data["document_text"]
        metadata = result.data["metadata"]
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gagal mengambil data sesi: {str(e)}")

    try:
        history_result = supabase.table("chat_history") \
            .select("role, content") \
            .eq("session_id", body.sessionId) \
            .order("created_at") \
            .limit(20) \
            .execute()
        history = history_result.data or []
    except Exception:
        history = []

    MAX_CHARS = 40000
    if len(doc_text) <= MAX_CHARS:
        truncated_doc = doc_text
    else:
        head = doc_text[:28000]
        tail = doc_text[-12000:]
        truncated_doc = head + "\n\n...[bagian tengah dokumen diringkas karena panjang]...\n\n" + tail
        logger.info(f"Dokumen dipotong: {len(doc_text)} → {len(truncated_doc)} karakter (smart truncation)")

    if body.isQuiz:
        qt = body.quizType
        if qt == 'essay':
            type_desc = "ESSAY — soal uraian yang membutuhkan analisis. Tidak ada pilihan jawaban."
            q_schema = '"type": "essay", "question": "pertanyaan essay analisis mendalam", "options": [], "answer": null, "explanation": null'
            extra_rule = 'Untuk essay: "options" wajib [], "answer" dan "explanation" wajib null.'
        elif qt == 'keduanya':
            type_desc = "campuran PILIHAN GANDA dan ESSAY (sekitar 50:50)."
            q_schema = '"type": "pilgan", "question": "...", "options": ["A. ...", "B. ...", "C. ...", "D. ..."], "answer": "A", "explanation": "..."'
            extra_rule = 'Buat campuran: pilgan (type:"pilgan", options 4 pilihan, answer A/B/C/D, explanation) dan essay (type:"essay", options:[], answer:null, explanation:null).'
        else:
            type_desc = "PILIHAN GANDA dengan 4 opsi (A/B/C/D)."
            q_schema = '"type": "pilgan", "question": "teks pertanyaan lengkap", "options": ["A. opsi A", "B. opsi B", "C. opsi C", "D. opsi D"], "answer": "A", "explanation": "penjelasan singkat"'
            extra_rule = 'PENTING: field "answer" diisi HANYA huruf kapital (A / B / C / D).'

        system_prompt = f"""Kamu adalah AI pembuat soal kuis. Buat soal {type_desc}
{extra_rule}


PENTING: Kembalikan HANYA JSON valid berikut, tanpa teks apapun di luar JSON:
{{
  "output": "quiz",
  "quiz": {{
    "questions": [
      {{{q_schema}}}
    ]
  }},
  "citations": []
}}

TEKS DOKUMEN (sumber soal):
---
{truncated_doc}
---"""
    else:
        system_prompt = f"""Kamu adalah asisten AI untuk membaca dan menganalisis dokumen.
PENTING: Kembalikan HANYA JSON berikut, tanpa teks lain di luar JSON:
{{"output": "jawaban lengkap di sini", "citations": [{{"page": nomor_halaman, "paragraph": nomor_paragraf}}]}}
Jika tidak ada kutipan spesifik, gunakan citations: []

ATURAN WAJIB:
1. Jawab HANYA berdasarkan teks dokumen — DILARANG menjawab dari pengetahuan umum di luar dokumen
2. DILARANG mengatakan tidak bisa membuka/membaca/mengakses dokumen
3. Untuk perintah terhadap dokumen (ringkas, terjemahkan, jelaskan isi) — langsung kerjakan
4. Jika pertanyaan membahas topik yang TIDAK ADA dalam dokumen, jawab: "Pertanyaan ini tidak dapat dijawab karena topik tersebut tidak dibahas dalam dokumen yang diupload."
5. Jawab dalam bahasa yang sama dengan pertanyaan user
6. FORMAT MARKDOWN dalam field "output":
   - Daftar bernomor: gunakan `1.`, `2.`, `3.` — satu item per baris, BUKAN inline
   - Daftar tanpa urutan: gunakan `- item` — BUKAN huruf a, b, c atau inline
   - Penekanan: gunakan **bold**
   - JANGAN tulis daftar dalam satu paragraf seperti "1) item, 2) item, 3) item"

Metadata dokumen:
- Judul: {metadata.get('name', '')}
- Topik: {metadata.get('topic', '')}
- Penulis: {metadata.get('author', '')}
- Tahun: {metadata.get('year', '')}

TEKS DOKUMEN:
---
{truncated_doc}
---"""

    messages = [{"role": "system", "content": system_prompt}]

    BAD_PATTERNS = ["tidak dapat membuka", "tidak bisa membuka", "tidak dapat mengakses",
                    "tidak bisa mengakses", "cannot open", "cannot access", "can't open",
                    "unable to access", "unable to open"]

    for h in history[-6:]:
        if h["role"] == "quiz_result":
            continue
        content = h["content"] or ""
        if h["role"] == "assistant" and any(p in content.lower() for p in BAD_PATTERNS):
            continue
        messages.append({"role": h["role"], "content": content})

    messages.append({"role": "user", "content": body.chatInput})

    max_tokens = 2500 if body.isQuiz else 1000

    try:
        response = groq_create(
            messages=messages,
            temperature=0.4,
            max_tokens=max_tokens,
            response_format={"type": "json_object"} if body.isQuiz else None,
            allow_fallback=not body.isQuiz,  # quiz butuh model penuh, jangan fallback
        )
        raw = response.choices[0].message.content.strip()
        logger.info(f"AI raw response (isQuiz={body.isQuiz}): {raw[:300]}")
    except Exception as e:
        err = str(e)
        if "rate_limit_exceeded" in err:
            raise HTTPException(status_code=429, detail=parse_rate_limit_error(err))
        raise HTTPException(status_code=500, detail=f"Gagal memanggil AI: {err}")

    try:
        ai_data = parse_json_from_ai(raw)
        output = ai_data.get("output", raw)
        quiz = ai_data.get("quiz", None)
        citations = ai_data.get("citations", [])
        if body.isQuiz:
            logger.info(f"Quiz parsed: {len(quiz.get('questions', [])) if quiz else 0} soal")
    except (json.JSONDecodeError, ValueError):
        logger.warning(f"Gagal parse JSON dari AI. Raw: {raw[:200]}")
        output = raw
        quiz = None
        citations = []

    try:
        supabase.table("chat_history").insert(
            {"session_id": body.sessionId, "role": "user", "content": body.chatInput, "user_id": body.userId}
        ).execute()
        supabase.table("chat_history").insert(
            {"session_id": body.sessionId, "role": "assistant", "content": output, "user_id": body.userId}
        ).execute()
    except Exception:
        pass

    return {"output": output, "quiz": quiz, "citations": citations}


class QuizResultRequest(BaseModel):
    session_id: str
    user_id: str | None = None
    total_score: int
    pilgan_correct: int | None = None
    pilgan_total: int | None = None
    essay_avg: int | None = None
    essay_total: int | None = None


@app.post("/save-quiz-result")
async def save_quiz_result(body: QuizResultRequest):
    content = json.dumps({
        "total_score": body.total_score,
        "pilgan_correct": body.pilgan_correct,
        "pilgan_total": body.pilgan_total,
        "essay_avg": body.essay_avg,
        "essay_total": body.essay_total,
    })
    try:
        supabase.table("chat_history").insert({
            "session_id": body.session_id,
            "role": "quiz_result",
            "content": content,
            "user_id": body.user_id,
        }).execute()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gagal menyimpan hasil quiz: {str(e)}")
    return {"ok": True}


class EssayGradeRequest(BaseModel):
    session_id: str
    question: str
    user_answer: str


@app.post("/grade-essay")
async def grade_essay(body: EssayGradeRequest):
    try:
        result = supabase.table("sessions") \
            .select("document_text") \
            .eq("session_id", body.session_id) \
            .single() \
            .execute()
        if not result.data:
            raise HTTPException(status_code=404, detail="Sesi tidak ditemukan.")
        doc_text = result.data["document_text"][:8000]
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gagal mengambil data sesi: {str(e)}")

    messages = [
        {
            "role": "system",
            "content": f"""Kamu adalah penilai jawaban essay berdasarkan dokumen.

TEKS DOKUMEN (konteks penilaian):
---
{doc_text}
---

Nilai jawaban user untuk soal essay berikut. Kembalikan HANYA JSON valid:
{{"score": <angka 0-100>, "feedback": "umpan balik spesifik tentang jawaban user", "ideal_answer": "jawaban ideal berdasarkan dokumen"}}

Kriteria:
- 90-100: Sangat lengkap dan akurat
- 70-89: Mencakup poin utama
- 50-69: Sebagian benar tapi kurang lengkap
- 0-49: Tidak tepat atau sangat kurang""",
        },
        {
            "role": "user",
            "content": f"Soal: {body.question}\n\nJawaban user: {body.user_answer}",
        },
    ]

    try:
        response = groq_create(
            messages=messages,
            temperature=0.3,
            max_tokens=600,
            response_format={"type": "json_object"},
        )
        raw = response.choices[0].message.content.strip()
        data = parse_json_from_ai(raw)
        return {
            "score": max(0, min(100, int(data.get("score", 0)))),
            "feedback": data.get("feedback", ""),
            "ideal_answer": data.get("ideal_answer", ""),
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gagal menilai essay: {str(e)}")
